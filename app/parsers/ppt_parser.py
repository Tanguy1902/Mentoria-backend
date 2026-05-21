"""
PowerPoint text extraction using python-pptx.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches  # noqa: F401 – imported for type completeness

from app.core.exceptions import ParsingError
from app.core.logging import get_logger
from app.parsers.base import BaseParser, ParsedDocument

logger = get_logger(__name__)

SUPPORTED_EXTENSIONS = {".pptx", ".ppt"}


class PPTParser(BaseParser):
    """Extracts text from PowerPoint presentations."""

    def supports(self, file_extension: str) -> bool:
        return file_extension.lower() in SUPPORTED_EXTENSIONS

    def parse(self, file_path: Path) -> ParsedDocument:
        """Extract text from all slides and shapes in a PowerPoint file."""
        logger.info("Parsing PowerPoint: %s", file_path.name)

        if not file_path.exists():
            raise ParsingError(f"PowerPoint file not found: {file_path}")

        try:
            prs = Presentation(str(file_path))
        except Exception as exc:
            raise ParsingError(
                f"Failed to open PowerPoint: {file_path.name}",
                detail=str(exc),
            ) from exc

        slides_text: list[str] = []
        slide_count = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []

            for shape in slide.shapes:
                # Extract text from text frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = paragraph.text.strip()
                        if para_text:
                            slide_parts.append(para_text)

                # Extract text from tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = [
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text.strip()
                        ]
                        if row_texts:
                            slide_parts.append(" | ".join(row_texts))

            if slide_parts:
                slide_text = f"--- Slide {slide_num} ---\n" + "\n".join(slide_parts)
                slides_text.append(slide_text)

        full_text = "\n\n".join(slides_text)

        if not full_text.strip():
            raise ParsingError(
                f"No extractable text found in PowerPoint: {file_path.name}",
                detail="The presentation may contain only images or embedded objects.",
            )

        logger.info(
            "PowerPoint parsed successfully: %d slides, %d characters",
            slide_count,
            len(full_text),
        )

        return ParsedDocument(
            text=full_text,
            page_count=slide_count,
            metadata={
                "parser": "python-pptx",
                "filename": file_path.name,
            },
        )
